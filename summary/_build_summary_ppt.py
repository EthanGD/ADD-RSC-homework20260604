import json
import os
from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


FONT_NAME = "微软雅黑"


def _set_run_font(run, size_pt, bold=False, color_rgb=None):
    run.font.name = FONT_NAME
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color_rgb is not None:
        run.font.color.rgb = color_rgb


def _style_text_frame(tf, size_pt, bold=False, color_rgb=None):
    for p in tf.paragraphs:
        for r in p.runs:
            _set_run_font(r, size_pt=size_pt, bold=bold, color_rgb=color_rgb)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    _style_text_frame(slide.shapes.title.text_frame, size_pt=40, bold=True)
    sub = slide.placeholders[1]
    sub.text = subtitle
    _style_text_frame(sub.text_frame, size_pt=18, bold=False, color_rgb=RGBColor(0x44, 0x44, 0x44))
    return slide


def add_bullets_slide(prs, title, bullets, font_size=18):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    _style_text_frame(slide.shapes.title.text_frame, size_pt=30, bold=True)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    first = True
    for item in bullets:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        for r in p.runs:
            _set_run_font(r, size_pt=font_size, bold=False)

    return slide


def add_picture_slide(prs, title, image_path, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    _style_text_frame(slide.shapes.title.text_frame, size_pt=30, bold=True)

    left = Inches(0.7)
    top = Inches(1.5)
    width = Inches(12.0)

    if os.path.exists(image_path):
        pic = slide.shapes.add_picture(image_path, left, top, width=width)
        if caption:
            tx = slide.shapes.add_textbox(Inches(0.7), pic.top + pic.height + Inches(0.1), Inches(12.0), Inches(0.6))
            tx.text_frame.text = caption
            _style_text_frame(tx.text_frame, size_pt=14, bold=False, color_rgb=RGBColor(0x55, 0x55, 0x55))
    else:
        tx = slide.shapes.add_textbox(Inches(0.7), top, Inches(12.0), Inches(1.0))
        tx.text_frame.text = "找不到图片：" + image_path
        _style_text_frame(tx.text_frame, size_pt=16, bold=False, color_rgb=RGBColor(0xAA, 0x00, 0x00))

    return slide


def safe_read_json(path):
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return None


def walk_files(root_dir):
    files = []
    for dirpath, dirnames, filenames in os.walk(root_dir):
        dirnames[:] = [d for d in dirnames if d not in {"__pycache__", ".git", ".venv", "venv"}]
        for fn in filenames:
            if fn.endswith((".pyc", ".pth")):
                continue
            files.append(os.path.join(dirpath, fn))
    return files


def main():
    root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    out_pptx = os.path.join(os.path.dirname(__file__), "ADD4RSC_复现项目总结.pptx")

    train_args_path = os.path.join(root, "resnet50Save", "save", "icbhi_resnet50_train_resnet", "train_args.json")
    results_path = os.path.join(root, "resnet50Save", "save", "icbhi_resnet50_train_resnet", "results.json")
    train_args = safe_read_json(train_args_path) or {}
    results = safe_read_json(results_path) or {}

    best_pth = os.path.join(root, "resnet50Save", "save", "icbhi_resnet50_train_resnet", "best.pth")
    best_epoch_pth = os.path.join(root, "resnet50Save", "save", "icbhi_resnet50_train_resnet", "best_epoch_15.pth")

    hp_lines = []
    for k, label in [
        ("model", "模型"),
        ("epochs", "epochs"),
        ("batch_size", "batch_size"),
        ("learning_rate", "lr"),
        ("weight_decay", "weight_decay"),
        ("optimizer", "optimizer"),
        ("sample_rate", "sample_rate"),
        ("desired_length", "cycle秒"),
        ("n_mels", "n_mels"),
        ("nfft", "nfft"),
        ("specaug_policy", "SpecAug"),
        ("loss_beta", "loss_beta"),
        ("denoise_d_model", "denoise_d_model"),
        ("denoise_num_heads", "denoise_heads"),
        ("denoise_depth", "denoise_depth"),
        ("ma_update", "EMA"),
        ("ma_beta", "EMA_beta"),
        ("data_folder", "data_folder"),
    ]:
        if k in train_args:
            hp_lines.append(f"{label}: {train_args[k]}")

    result_lines = []
    if results:
        for exp, v in results.items():
            if isinstance(v, list) and len(v) == 3:
                result_lines.append(f"{exp}: Sp={v[0]}  Se={v[1]}  Score={v[2]}")
            else:
                result_lines.append(f"{exp}: {v}")
    else:
        result_lines.append("未找到 results.json 或解析失败")

    structure = [
        "main.py：训练入口（ResNet50/AST + ADD 去噪模块 + 指标计算/保存）",
        "models/：ResNet50、AST、ADD 相关模块（AFF/MHDA/DDL）与损失",
        "util/：ICBHI 数据集处理、特征提取、增强、训练工具函数",
        "download_*.py：下载 ICBHI 数据与 AST 预训练权重的脚本",
        "predict.py：推理脚本骨架（目前单文件预测仍需补全预处理）",
        "resnet50Save/：当前已跑出的 ResNet50 训练产物（pth、曲线、json）",
        "大量 .md/.pdf/.pptx：论文/模块理解笔记与演示材料",
    ]

    issues = [
        "ICBHIDataset 使用随机 60/40 split（seed=1），未真正实现 official fold（util/icbhi_dataset.py）。",
        "validate() 的 loss 口径与训练不一致（train 用 loss_beta 加权；validate 直接相加）（main.py）。",
        "MHDA 内部 LayerNorm 固定到 cuda:0（models/adapt_diff_denoise.py），可能导致设备不兼容。",
        "predict.py 缺少与训练一致的音频→fbank 预处理。",
        "AST 强依赖 timm==0.4.5（models/ast.py assert）。",
    ]

    next_steps = [
        "对齐数据划分（官方 split / RespireNet 80-20），再对齐论文超参。",
        "补全 AST backbone 的训练与评测闭环。",
        "处理类别不平衡：--weighted_sampler / --weighted_loss，重点观察 Se/Score。",
        "修复 MHDA 设备硬编码与验证口径，提升可复现性。",
    ]

    all_files = walk_files(root)
    py_files = [p for p in all_files if p.lower().endswith(".py")]
    md_files = [p for p in all_files if p.lower().endswith(".md")]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    today = date.today().isoformat()
    add_title_slide(
        prs,
        "ADD4RSC 复现项目总结（本仓库）",
        "原始仓库：github.com/deegy666/ADD-RSC\n本地路径：" + root + "\n生成日期：" + today,
    )

    add_bullets_slide(
        prs,
        "项目目标与现状",
        [
            "目标：复现 Interspeech 2025 论文《Adaptive Differential Denoising for Respiratory Sounds Classification》的训练/评测流程。",
            "当前进度：已完成 ResNet50 backbone 的训练跑通；产物位于 resnet50Save/。",
            "尚未覆盖：AST backbone 完整复现；对齐论文 official 数据划分与最终指标。",
        ],
        font_size=18,
    )

    add_bullets_slide(prs, "仓库结构一览", structure, font_size=18)

    add_bullets_slide(
        prs,
        "核心训练流水线（main.py）",
        [
            "数据：ICBHI wav+txt → 按呼吸周期切分 → 8s 定长/补齐",
            "特征：kaldi fbank(n_mels=128) → Resize(1024,256) → SpecAugment(训练)",
            "去噪：DiffTransformerLayer（AFF(AFNO1D)+MHDA+SwiGLU）",
            "主干：ResNet50(单通道) 或 AST(ViT) → 线性分类器",
            "指标：Sp/Se/Score；保存：best.pth + train_args.json + 曲线图",
        ],
        font_size=18,
    )

    add_picture_slide(prs, "论文整体架构示意", os.path.join(root, "image", "fig_0216.png"), "来自仓库 README")

    add_bullets_slide(prs, "ResNet50 训练超参（train_args.json）", hp_lines or ["未读取到 train_args.json"], font_size=16)

    add_picture_slide(
        prs,
        "训练过程曲线（checkpoints.png）",
        os.path.join(root, "resnet50Save", "save", "icbhi_resnet50_train_resnet", "checkpoints.png"),
    )

    add_bullets_slide(
        prs,
        "当前结果与产物位置",
        [
            "结果文件：" + results_path,
            *result_lines,
            "模型权重：" + best_pth,
            "备份权重：" + best_epoch_pth,
        ],
        font_size=16,
    )

    add_bullets_slide(prs, "对齐论文的风险点", issues, font_size=16)
    add_bullets_slide(prs, "下一步建议", next_steps, font_size=18)

    add_bullets_slide(
        prs,
        "附录：文件统计",
        [
            "总文件数（排除 .pyc/.pth 与 __pycache__）：{}".format(len(all_files)),
            "Python 源码：{}".format(len(py_files)),
            "Markdown 文档：{}".format(len(md_files)),
        ],
        font_size=18,
    )

    prs.save(out_pptx)
    print(out_pptx)


if __name__ == "__main__":
    main()

