#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""生成ADD4RSC论文中文PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============ 颜色定义 ============
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # 深蓝紫背景
BG_GRADIENT = RGBColor(0x16, 0x21, 0x3E)    # 渐变蓝
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)    # 亮蓝强调色
ACCENT_GREEN = RGBColor(0x4E, 0xCB, 0x71)   # 绿色强调色
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x42)  # 橙色强调色
ACCENT_PURPLE = RGBColor(0xA0, 0x5C, 0xF0)  # 紫色强调色
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
TEXT_DIM = RGBColor(0x88, 0x88, 0x99)
HIGHLIGHT_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
CARD_BG = RGBColor(0x22, 0x2B, 0x45)

def add_bg(slide, color=BG_DARK):
    """为幻灯片添加背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    """添加矩形装饰块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_para(slide, left, top, width, height, paragraphs_data, font_name='微软雅黑'):
    """添加多段落文本框
    paragraphs_data: list of (text, font_size, color, bold, alignment)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, font_size, color, bold, alignment) in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
    return txBox

def add_card(slide, left, top, width, height, title, content_lines, accent_color=ACCENT_BLUE, title_size=20, content_size=14):
    """添加信息卡片"""
    # 卡片背景
    card = add_rect(slide, left, top, width, height, CARD_BG)
    # 左侧强调线
    add_rect(slide, left, top, Inches(0.08), height, accent_color)
    # 标题
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.3), Inches(0.4),
                 title, title_size, accent_color, True)
    # 内容
    para_data = []
    for line in content_lines:
        para_data.append((line, content_size, TEXT_LIGHT, False, PP_ALIGN.LEFT))
    if para_data:
        add_multi_para(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.3), height - Inches(0.7),
                       para_data)

# ================================================================
# Slide 1: 标题页
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_DARK)
# 顶部装饰线
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_BLUE)
# 左侧装饰块
add_rect(slide, Inches(0), Inches(0), Inches(0.06), prs.slide_height, ACCENT_PURPLE)

# 标题
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "自适应差分去噪网络", 44, TEXT_WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
             "面向呼吸音分类的自适应差分去噪方法", 30, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.6),
             "Adaptive Differential Denoising for Respiratory Sounds Classification", 20, TEXT_DIM, False, PP_ALIGN.CENTER)

# 作者与会议信息
add_text_box(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.4),
             "Gaoyang Dong · Zhicheng Zhang · Ping Sun · Minghui Zhang", 18, TEXT_LIGHT, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.4),
             "南昌大学信息工程学院 · 复旦大学计算机科学与技术学院", 16, TEXT_DIM, False, PP_ALIGN.CENTER)

# 会议标签
conf_box = add_rect(slide, Inches(5.2), Inches(5.8), Inches(2.8), Inches(0.5), ACCENT_BLUE)
add_text_box(slide, Inches(5.2), Inches(5.8), Inches(2.8), Inches(0.5),
             "Interspeech 2025", 22, TEXT_WHITE, True, PP_ALIGN.CENTER)

# 性能亮点
add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
             "ICBHI 2017 Score: 65.53%  |  超越SOTA +1.99%", 20, HIGHLIGHT_YELLOW, True, PP_ALIGN.CENTER)

# ================================================================
# Slide 2: 研究背景与动机
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.5),
             "01 研究背景与动机", 28, ACCENT_BLUE, True)

# 左侧 - 问题挑战
add_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(3.0),
         "核心挑战", [
             "■ 自动呼吸音分类(ARSC)面临两大难题:",
             "   1. 真实场景背景噪声干扰(心音、环境噪声等)",
             "   2. 标注样本稀缺，数据不足",
             "",
             "■ 传统方法缺陷:",
             "   · 谱减法/维纳滤波：需纯净参考信号",
             "   · Butterworth滤波：固定截止频率不自适应",
             "   · 深度去噪模型：需配对干净/噪声数据",
             "   · 噪声与诊断特征频谱重叠(如心音掩盖细裂纹音)"
         ], ACCENT_ORANGE, 18, 14)

# 右侧 - 研究意义
add_card(slide, Inches(6.5), Inches(1.2), Inches(6.3), Inches(3.0),
         "研究意义", [
             "■ 非侵入式肺部疾病诊断工具",
             "   · 哮喘、肺炎、COPD等疾病筛查",
             "",
             "■ 远程医疗应用价值",
             "   · 资源匮乏地区早期筛查",
             "   · 数字听诊器与可穿戴传感器",
             "",
             "■ 首次探索深度学习去噪技术用于呼吸音",
             "   · 摆脱'先去噪再分类'的传统顺序范式",
             "   · 隐式去噪，无需干净标签"
         ], ACCENT_GREEN, 18, 14)

# 底部 - 创新点概览
add_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.06), ACCENT_BLUE)
add_text_box(slide, Inches(0.5), Inches(4.7), Inches(12), Inches(0.4),
             "三大核心创新", 22, HIGHLIGHT_YELLOW, True, PP_ALIGN.CENTER)

# 三列创新卡片
col_w = Inches(3.9)
col_gap = Inches(0.3)
innovations = [
    ("自适应频域滤波器(AFF)", ACCENT_BLUE,
     ["· 可学习谱掩码 + 软阈值(SoftShrink)",
      "· 消除噪声同时保留诊断性高频成分",
      "· 自适应实例级掩码机制"]),
    ("差分去噪层(DDL)", ACCENT_GREEN,
     ["· 基于差分注意力机制(MHDA)",
      "· 通过增强样本对比抑制噪声扰动",
      "· 双路注意力差分: σ(Q₁K₁ᵀ) - λ·σ(Q₂K₂ᵀ)"]),
    ("偏置去噪损失(BDL)", ACCENT_PURPLE,
     ["· 联合优化分类精度与鲁棒性",
      "· 无需纯净标签/配对数据",
      "· 标签平滑注入医学先验知识"]),
]
for i, (title, color, lines) in enumerate(innovations):
    left = Inches(0.5) + i * (col_w + col_gap)
    add_card(slide, left, Inches(5.2), col_w, Inches(1.8), title, lines, color, 16, 13)

# ================================================================
# Slide 3: 模型架构总览
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.5),
             "02 模型架构总览", 28, ACCENT_BLUE, True)

# 架构流程 - 横向展示
flow_steps = [
    ("预处理", "重采样16kHz\n8秒定长\nMel频谱图"),
    ("AFF", "自适应频域\n滤波器\nFFT→Linear→ReLU\n→Linear→SoftShrink\n→IFFT"),
    ("DDL", "差分去噪层\nLN→MHDA→Add\n→LN→SwiGLU\n→Add"),
    ("Backbone", "主干网络\nResNet50 或\nAST-Base"),
    ("Classifier", "分类器\n4类输出"),
]
step_w = Inches(2.2)
step_h = Inches(2.8)
for i, (name, detail) in enumerate(flow_steps):
    left = Inches(0.5) + i * Inches(2.45)
    top = Inches(1.3)
    color = [ACCENT_ORANGE, ACCENT_BLUE, ACCENT_GREEN, RGBColor(0x55, 0x5E, 0x7B), ACCENT_PURPLE][i]
    # 卡片
    card = add_rect(slide, left, top, step_w, step_h, CARD_BG)
    add_rect(slide, left, top, step_w, Inches(0.06), color)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.15), step_w - Inches(0.2), Inches(0.4),
                 name, 18, color, True, PP_ALIGN.CENTER)
    # 详情
    para_data = [(line, 12, TEXT_LIGHT, False, PP_ALIGN.CENTER) for line in detail.split('\n')]
    add_multi_para(slide, left + Inches(0.1), top + Inches(0.6), step_w - Inches(0.2), step_h - Inches(0.7),
                   para_data)
    # 箭头连接
    if i < len(flow_steps) - 1:
        arrow_left = left + step_w
        add_text_box(slide, arrow_left, top + Inches(1.3), Inches(0.25), Inches(0.3),
                     "→", 20, ACCENT_BLUE, True, PP_ALIGN.CENTER)

# 损失函数
add_rect(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.5), CARD_BG)
add_rect(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.06), HIGHLIGHT_YELLOW)
add_text_box(slide, Inches(0.7), Inches(4.45), Inches(11), Inches(0.4),
             "联合损失函数", 22, HIGHLIGHT_YELLOW, True)

add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.6),
             "ℒ = β · ℒ_Bias Denoise + (1-β) · ℒ_CE", 24, TEXT_WHITE, True, PP_ALIGN.CENTER)

para_data = [
    ("ℒ_Bias Denoise: 偏置去噪损失 — 标签平滑交叉熵, 不需纯净标签, 注入医学先验知识", 14, TEXT_LIGHT, False, PP_ALIGN.LEFT),
    ("ℒ_CE: 分类交叉熵损失 — 保障分类精度", 14, TEXT_LIGHT, False, PP_ALIGN.LEFT),
    ("β = 0.5: 控制去噪引导与分类精度之间的平衡", 14, TEXT_LIGHT, False, PP_ALIGN.LEFT),
]
add_multi_para(slide, Inches(0.7), Inches(5.7), Inches(11), Inches(1.5), para_data)

# ================================================================
# Slide 4: 自适应频域滤波器(AFF)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(5), Inches(0.5),
             "03 自适应频域滤波器 (AFF)", 28, ACCENT_BLUE, True)

# 左侧 - 算法原理
add_card(slide, Inches(0.5), Inches(1.2), Inches(6.3), Inches(3.5),
         "算法原理", [
             "Step 1: FFT变换 → 获取频域表示 X(U,V)",
             "   X(U,V) = Σ Σ X(t,f) · e^(-2πi(Ut+Vf))",
             "",
             "Step 2: 可学习谱掩码 ℳ(·) 滤波",
             "   由 Linear → ReLU → Linear 组成",
             "   实例自适应，逐样本生成掩码",
             "",
             "Step 3: 软阈值(SoftShrink)稀疏化",
             "   Sα(x) = sign(x) · max{|x| - α, 0}",
             "   α=0.02 控制稀疏度，抑制噪声",
             "",
             "Step 4: IFFT逆变换 → 回到时频域",
             "   X̂(t,f) = F⁻¹[Sα(ℳ(X(U,V))) ⊙ X(U,V)]",
         ], ACCENT_BLUE, 17, 14)

# 右侧 - 设计特点
add_card(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(2.0),
         "核心设计特点", [
             "★ 频域卷积定理: 时域卷积 = 频域逐元素乘",
             "★ 可学习实例自适应掩码 ℳ(·)",
             "★ SoftShrink软阈值 → 噪声抑制+稀疏化",
             "★ 残差连接: 保留原始信号信息"
         ], ACCENT_ORANGE, 17, 14)

# 右侧 - 优势对比
add_card(slide, Inches(7.0), Inches(3.4), Inches(5.8), Inches(1.3),
         "vs 传统滤波方法", [
             "传统方法: 固定截止频率, 非自适应, 需纯净参考",
             "AFF方法: 可学习掩码, 实例自适应, 无需纯净数据"
         ], ACCENT_GREEN, 17, 14)

# 公式框
add_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.8), CARD_BG)
add_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.06), ACCENT_BLUE)
add_text_box(slide, Inches(0.7), Inches(5.15), Inches(11), Inches(0.35),
             "关键公式", 20, ACCENT_BLUE, True)
add_text_box(slide, Inches(1.5), Inches(5.55), Inches(10), Inches(0.4),
             "X̂(t,f) = F⁻¹ [ Sα(ℳ(X(U,V))) ⊙ X(U,V) ]", 22, TEXT_WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.4),
             "Sα(x) = sign(x) · max{|x| - α, 0}      α = 0.02", 18, TEXT_LIGHT, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.4), Inches(10), Inches(0.4),
             "ℳ(·): 可学习实例自适应掩码 = Linear → ReLU → Linear", 16, TEXT_DIM, False, PP_ALIGN.CENTER)

# ================================================================
# Slide 5: 差分去噪层(DDL)与MHDA
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(6), Inches(0.5),
             "04 差分去噪层 (DDL) 与多头差分注意力 (MHDA)", 28, ACCENT_GREEN, True)

# DDL流程
add_card(slide, Inches(0.5), Inches(1.2), Inches(6.3), Inches(2.8),
         "差分去噪层(DDL) 流程", [
             "Input → LayerNorm → MHDA → 残差加法 → LayerNorm → SwiGLU → 残差加法",
             "",
             "■ LN(层归一化): 稳定特征分布",
             "   减轻不同录音条件带来的变化",
             "",
             "■ MHDA(多头差分注意力): 核心创新",
             "   双路注意力差分，抑制噪声敏感变化",
             "   保留稳定、高置信度的分类关键特征",
             "",
             "■ SwiGLU: 非线性激活 + 平滑梯度流",
             "   保留判别性特征，确保信息不丢失"
         ], ACCENT_GREEN, 17, 14)

# MHDA详细机制
add_card(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(2.8),
         "MHDA 差分注意力机制", [
             "投影分解:",
             "  [Q₁; Q₂] = X·Wᵠ,  [K₁; K₂] = X·Wᵏ,  V = X·Wᵛ",
             "",
             "差分注意力计算:",
             "  A = (σ(Q₁K₁ᵀ/√d) - λ·σ(Q₂K₂ᵀ/√d)) · V",
             "",
             "λ: 可学习缩放因子",
             "  λ = exp(q₁·k₁) - exp(q₂·k₂) + λ_init",
             "  λ_init = 0.8 - 0.6·exp(-0.3·depth)",
             "",
             "后处理: GroupNorm → Concat → Linear"
         ], ACCENT_PURPLE, 17, 14)

# 核心思想对比
add_rect(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.8), CARD_BG)
add_rect(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.06), ACCENT_GREEN)
add_text_box(slide, Inches(0.7), Inches(4.35), Inches(11), Inches(0.35),
             "差分注意力的核心思想", 20, ACCENT_GREEN, True)

# 左右对比
add_card(slide, Inches(0.7), Inches(4.8), Inches(5.8), Inches(2.0),
         "传统注意力", [
             "σ(QKᵀ) · V",
             "· 单路softmax注意力",
             "· 对噪声敏感，无法区分噪声与信号",
             "· 噪声变化会直接影响注意力分布"
         ], RGBColor(0x88, 0x44, 0x44), 17, 14)

add_card(slide, Inches(6.7), Inches(4.8), Inches(5.8), Inches(2.0),
         "差分注意力 (MHDA)", [
             "(σ(Q₁K₁ᵀ) - λ·σ(Q₂K₂ᵀ)) · V",
             "· 双路注意力差分",
             "· 过滤噪声敏感变化",
             "· 保留稳定高置信度特征"
         ], ACCENT_GREEN, 17, 14)

# ================================================================
# Slide 6: 偏置去噪损失(BDL)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_PURPLE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(5), Inches(0.5),
             "05 偏置去噪损失 (Bias Denoise Loss)", 28, ACCENT_PURPLE, True)

# 左侧 - 损失函数组成
add_card(slide, Inches(0.5), Inches(1.2), Inches(6.3), Inches(3.0),
         "联合损失函数", [
             "ℒ = β · ℒ_Bias Denoise + (1-β) · ℒ_CE",
             "",
             "■ ℒ_CE: 标准交叉熵分类损失",
             "   保障分类精度的基础损失",
             "",
             "■ ℒ_Bias Denoise: 偏置去噪损失",
             "   ℒ_BD = -Σ [y_c(1-ε) + ε/C] · log[φ(Norm(p))]",
             "   · ε=0.2: 标签平滑参数",
             "   · C=4: 分类类别数",
             "   · φ: 1×1卷积映射去噪特征",
             "   · p: DDL输出的去噪特征",
             "",
             "■ β=0.5: 平衡去噪与分类的权重"
         ], ACCENT_PURPLE, 17, 14)

# 右侧 - 设计思想
add_card(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(3.0),
         "设计思想与优势", [
             "★ 无需纯净标签/配对数据",
             "   只依赖类别标签 — 适合真实部署场景",
             "",
             "★ 标签平滑注入医学先验",
             "   ε作为不确定性缓冲",
             "   防止过度自信地移除模糊频段",
             "",
             "★ 防止诊断信息丢失",
             "   临床噪声与目标信号频谱重叠",
             "   (如心音掩盖细裂纹音)",
             "   ε防止错误删除含诊断信息的频段",
             "",
             "★ 联合优化: 分类精度 + 噪声鲁棒性"
         ], ACCENT_ORANGE, 17, 14)

# 公式展示
add_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), CARD_BG)
add_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.06), ACCENT_PURPLE)
add_text_box(slide, Inches(0.7), Inches(4.65), Inches(11), Inches(0.35),
             "关键公式", 20, ACCENT_PURPLE, True)

add_text_box(slide, Inches(1.5), Inches(5.1), Inches(10), Inches(0.4),
             "ℒ = β · ℒ_Bias Denoise + (1-β) · ℒ_CE", 24, TEXT_WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.5),
             "ℒ_BD = -Σ [y_c(1-ε) + ε/C] · log[φ(Norm(p))]", 20, ACCENT_PURPLE, False, PP_ALIGN.CENTER)

para_data = [
    ("β=0.5 | ε=0.2 | C=4 | α=0.02 | φ: 1×1 Conv", 16, TEXT_DIM, False, PP_ALIGN.CENTER),
    ("不需要配对干净/噪声数据 — 仅依赖类别标签 — 关键适用于真实场景部署", 14, TEXT_LIGHT, False, PP_ALIGN.CENTER),
]
add_multi_para(slide, Inches(0.7), Inches(6.2), Inches(11), Inches(1.0), para_data)

# ================================================================
# Slide 7: 实验结果 - SOTA对比
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(5), Inches(0.5),
             "06 实验结果 — SOTA对比", 28, ACCENT_BLUE, True)

# CNN方法表
add_text_box(slide, Inches(0.5), Inches(1.1), Inches(6), Inches(0.4),
             "CNN-based 方法对比", 18, ACCENT_ORANGE, True)

cnn_data = [
    ("LungRN+NL", "63.20", "41.32", "52.26"),
    ("RespireNet", "72.30", "40.10", "56.20"),
    ("StochNorm", "78.86", "36.40", "57.63"),
    ("CoTuning", "79.34", "37.24", "58.29"),
    ("SCL", "75.95", "39.15", "57.55"),
    ("Ours(ResNet50)", "83.76", "34.18", "58.97"),
]

# 表头
headers = ["模型", "Sp(%)", "Se(%)", "Score(%)"]
col_widths = [Inches(2.0), Inches(1.2), Inches(1.2), Inches(1.2)]
col_starts = [Inches(0.5), Inches(2.5), Inches(3.7), Inches(4.9)]
top_start = Inches(1.55)

for j, hdr in enumerate(headers):
    add_text_box(slide, col_starts[j], top_start, col_widths[j], Inches(0.3),
                 hdr, 13, ACCENT_BLUE, True, PP_ALIGN.CENTER)

for i, (model, sp, se, score) in enumerate(cnn_data):
    row_top = top_start + Inches(0.35) + i * Inches(0.3)
    color = HIGHLIGHT_YELLOW if "Ours" in model else TEXT_LIGHT
    bold = "Ours" in model
    vals = [model, sp, se, score]
    for j, val in enumerate(vals):
        add_text_box(slide, col_starts[j], row_top, col_widths[j], Inches(0.3),
                     val, 12, color, bold, PP_ALIGN.CENTER)

# Transformer方法表
add_text_box(slide, Inches(7.0), Inches(1.1), Inches(6), Inches(0.4),
             "Transformer-based 方法对比", 18, ACCENT_GREEN, True)

trans_data = [
    ("AFT on Mixed-500", "80.72", "42.86", "61.79"),
    ("AST Fine-tuning", "77.14", "41.97", "59.55"),
    ("Patch-Mix CL", "81.66", "43.07", "62.37"),
    ("M2D", "81.51", "45.08", "63.29"),
    ("DAT", "77.11", "42.50", "59.81"),
    ("BTS", "81.40", "45.67", "63.54"),
    ("LungAdapter", "80.43", "44.37", "62.40"),
    ("Ours(AST)", "85.13", "45.94", "65.53"),
]

col_starts2 = [Inches(7.0), Inches(9.0), Inches(10.2), Inches(11.4)]
col_widths2 = [Inches(2.0), Inches(1.2), Inches(1.2), Inches(1.2)]

for j, hdr in enumerate(headers):
    add_text_box(slide, col_starts2[j], top_start, col_widths2[j], Inches(0.3),
                 hdr, 13, ACCENT_GREEN, True, PP_ALIGN.CENTER)

for i, (model, sp, se, score) in enumerate(trans_data):
    row_top = top_start + Inches(0.35) + i * Inches(0.3)
    color = HIGHLIGHT_YELLOW if "Ours" in model else TEXT_LIGHT
    bold = "Ours" in model
    vals = [model, sp, se, score]
    for j, val in enumerate(vals):
        add_text_box(slide, col_starts2[j], row_top, col_widths2[j], Inches(0.3),
                     val, 12, color, bold, PP_ALIGN.CENTER)

# 性能亮点
add_rect(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.8), CARD_BG)
add_rect(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.06), HIGHLIGHT_YELLOW)
add_text_box(slide, Inches(0.7), Inches(4.35), Inches(11), Inches(0.35),
             "性能亮点", 20, HIGHLIGHT_YELLOW, True)

highlights = [
    ("🏆 Score: 65.53%", "超越前最佳BTS(63.54%) +1.99%", ACCENT_BLUE),
    ("🏆 Sp: 85.13%", "超越RepAugment(82.47%) +2.66%", ACCENT_GREEN),
    ("🏆 Se: 45.94%", "超越BTS(45.67%) +0.27%", ACCENT_ORANGE),
]
for i, (metric, detail, color) in enumerate(highlights):
    left = Inches(0.7) + i * Inches(4.1)
    add_text_box(slide, left, Inches(4.85), Inches(3.8), Inches(0.5),
                 metric, 24, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(5.4), Inches(3.8), Inches(0.4),
                 detail, 14, TEXT_LIGHT, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.7), Inches(5.9), Inches(11), Inches(0.6),
             "CNN和Transformer两种架构均取得领先 → 验证了自适应差分去噪机制的架构无关性", 16, TEXT_LIGHT, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.7), Inches(6.4), Inches(11), Inches(0.5),
             "Sp和Se同时提升(传统呈权衡关系) → 有效将信号模式与噪声干扰解耦", 14, TEXT_DIM, False, PP_ALIGN.CENTER)

# ================================================================
# Slide 8: 消融实验
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_GREEN)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(5), Inches(0.5),
             "07 消融实验", 28, ACCENT_GREEN, True)

add_text_box(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.35),
             "Backbone: AST | 数据集: ICBHI 2017", 16, TEXT_DIM, False)

# 消融实验表
ablation_data = [
    ("w AFF", "82.47", "44.47", "63.47", "+0.93", "仅自适应频域滤波"),
    ("w DDL", "83.73", "44.03", "63.88", "+1.34", "仅差分去噪层"),
    ("w AFF+DDL", "84.41", "44.78", "64.60", "+2.06", "两组件组合"),
    ("w AFF+ℒ_BD", "83.76", "44.62", "64.19", "+1.65", "AFF+偏置去噪损失"),
    ("w DDL+ℒ_BD", "84.54", "44.32", "64.43", "+1.89", "DDL+偏置去噪损失"),
    ("Ours (完整)", "85.13", "45.94", "65.53", "+2.99", "三组件完整模型"),
]

headers = ["配置", "Sp(%)", "Se(%)", "Score(%)", "Δ Score", "说明"]
col_starts = [Inches(0.5), Inches(3.0), Inches(4.2), Inches(5.4), Inches(6.6), Inches(7.8)]
col_widths_abl = [Inches(2.5), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2), Inches(4.5)]

# 表头
for j, hdr in enumerate(headers):
    add_text_box(slide, col_starts[j], Inches(1.4), col_widths_abl[j], Inches(0.3),
                 hdr, 14, ACCENT_GREEN, True, PP_ALIGN.CENTER)

# 基准行
add_rect(slide, Inches(0.5), Inches(1.75), Inches(12.3), Inches(0.3), RGBColor(0x33, 0x33, 0x44))
add_text_box(slide, Inches(0.5), Inches(1.75), Inches(2.5), Inches(0.3),
             "Baseline (AST)", 13, TEXT_DIM, False, PP_ALIGN.LEFT)
add_text_box(slide, Inches(5.4), Inches(1.75), Inches(1.2), Inches(0.3),
             "62.54", 13, TEXT_DIM, False, PP_ALIGN.CENTER)

# 数据行
for i, (config, sp, se, score, delta, desc) in enumerate(ablation_data):
    row_top = Inches(2.15) + i * Inches(0.35)
    is_ours = "完整" in config
    color = HIGHLIGHT_YELLOW if is_ours else TEXT_LIGHT
    bold = is_ours
    vals = [config, sp, se, score, delta, desc]
    for j, val in enumerate(vals):
        add_text_box(slide, col_starts[j], row_top, col_widths_abl[j], Inches(0.3),
                     val, 13, color, bold, PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# 消融结论
add_rect(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8), CARD_BG)
add_rect(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.06), ACCENT_GREEN)
add_text_box(slide, Inches(0.7), Inches(4.45), Inches(11), Inches(0.35),
             "消融实验结论", 20, ACCENT_GREEN, True)

conclusions = [
    ("AFF单独", "Score +0.93% → 自适应频域滤波有效抑制噪声", ACCENT_BLUE),
    ("DDL单独", "Score +1.34% → 差分注意力比AFF更强的去噪能力", ACCENT_GREEN),
    ("AFF+DDL", "Score +2.06% → 两组件互补, 协同增强", ACCENT_ORANGE),
    ("完整模型", "Score +2.99% → 三组件协同达到最优", HIGHLIGHT_YELLOW),
]
for i, (label, detail, color) in enumerate(conclusions):
    left = Inches(0.7) + i * Inches(3.1)
    add_text_box(slide, left, Inches(4.9), Inches(2.8), Inches(0.4),
                 label, 16, color, True)
    add_text_box(slide, left, Inches(5.35), Inches(2.8), Inches(0.6),
                 detail, 12, TEXT_LIGHT, False)

add_text_box(slide, Inches(0.7), Inches(5.9), Inches(11), Inches(0.5),
             "每个组件都至关重要，三者整合产生最优分类精度", 16, TEXT_LIGHT, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.7), Inches(6.4), Inches(11), Inches(0.5),
             "DDL在噪声抑制和特征精炼中作用更显著 | AFF+DDL组合比单独任一组件提升更多", 14, TEXT_DIM, False, PP_ALIGN.CENTER)

# ================================================================
# Slide 9: 实验设置
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_ORANGE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.5),
             "08 实验设置", 28, ACCENT_ORANGE, True)

# 数据集
add_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.8),
         "数据集: ICBHI 2017", [
             "■ 126位受试者真实录音",
             "■ 包含心音、环境噪声、传感器伪影",
             "■ 采样率: 4kHz ~ 44.1kHz → 统一重采样16kHz",
             "■ 呼吸周期时长: 0.2s ~ 16.2s → 固定8秒",
             "   · 短录音: 循环填充",
             "   · 长录音: 截断",
             "■ 4类分类: Normal / Crackle / Wheeze / Both",
             "■ 官方60-40分割",
             "■ Mel频谱图: 64 Mel滤波器, 窗口1024"
         ], ACCENT_ORANGE, 17, 14)

# 训练设置
add_card(slide, Inches(6.5), Inches(1.2), Inches(6.3), Inches(2.8),
         "训练配置", [
             "■ 框架: PyTorch 2.3.1",
             "■ GPU: Tesla V100",
             "■ 优化器: Adam",
             "■ 学习率: 5e-5",
             "■ 权重衰减: 0.1",
             "■ Batch Size: 8",
             "■ Epochs: 50",
             "",
             "■ Backbone: ResNet50 & AST-Base",
             "   (AudioSet预训练)"
         ], ACCENT_BLUE, 17, 14)

# 评估指标
add_card(slide, Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.5),
         "评估指标 (ICBHI官方)", [
             "■ Se (灵敏度): 正确识别异常声音的比例",
             "   Se = #正确识别异常 / #总异常",
             "",
             "■ Sp (特异度): 正确识别正常声音的比例",
             "   Sp = #正确识别正常 / #总正常",
             "",
             "■ Score = (Se + Sp) / 2",
             "   异常与正常识别率的平均值"
         ], ACCENT_GREEN, 17, 14)

# 超参数
add_card(slide, Inches(6.5), Inches(4.2), Inches(6.3), Inches(2.5),
         "关键超参数", [
             "■ β = 0.5        (去噪损失权重)",
             "■ α = 0.02       (SoftShrink阈值)",
             "■ ε = 0.2        (标签平滑参数)",
             "■ d_model = 256  (DDL隐藏维度)",
             "■ num_heads = 8  (MHDA注意力头数)",
             "■ depth = 6      (DDL层数)",
             "■ λ_init = 0.8 - 0.6·exp(-0.3·depth)"
         ], ACCENT_PURPLE, 17, 14)

# ================================================================
# Slide 10: 总结与展望
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), HIGHLIGHT_YELLOW)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.5),
             "09 总结与展望", 28, HIGHLIGHT_YELLOW, True)

# 三大创新回顾
add_text_box(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.4),
             "核心创新回顾", 22, ACCENT_BLUE, True, PP_ALIGN.CENTER)

innovations_review = [
    ("自适应频域滤波器(AFF)", ACCENT_BLUE,
     "可学习谱掩码+SoftShrink\n频域自适应噪声抑制\n保留诊断性高频成分"),
    ("差分去噪层(DDL)", ACCENT_GREEN,
     "多头差分注意力(MHDA)\nσ(Q₁K₁ᵀ)-λ·σ(Q₂K₂ᵀ)\n噪声变化过滤+特征保留"),
    ("偏置去噪损失(BDL)", ACCENT_PURPLE,
     "标签平滑交叉熵\n无需纯净标签\n医学先验+不确定性缓冲"),
]

for i, (title, color, detail) in enumerate(innovations_review):
    left = Inches(0.5) + i * Inches(4.2)
    card = add_rect(slide, left, Inches(1.6), Inches(3.9), Inches(2.0), CARD_BG)
    add_rect(slide, left, Inches(1.6), Inches(3.9), Inches(0.06), color)
    add_text_box(slide, left + Inches(0.1), Inches(1.7), Inches(3.7), Inches(0.35),
                 title, 17, color, True, PP_ALIGN.CENTER)
    para_data = [(line, 13, TEXT_LIGHT, False, PP_ALIGN.CENTER) for line in detail.split('\n')]
    add_multi_para(slide, left + Inches(0.1), Inches(2.1), Inches(3.7), Inches(1.4), para_data)

# 性能总结
add_rect(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(1.5), CARD_BG)
add_rect(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.06), HIGHLIGHT_YELLOW)
add_text_box(slide, Inches(0.7), Inches(3.95), Inches(11), Inches(0.35),
             "性能总结", 20, HIGHLIGHT_YELLOW, True)

perf_data = [
    ("Score: 65.53%", "超越SOTA +1.99%", HIGHLIGHT_YELLOW),
    ("Sp: 85.13%", "超越RepAugment +2.66%", ACCENT_GREEN),
    ("Se: 45.94%", "超越BTS +0.27%", ACCENT_ORANGE),
]
for i, (metric, detail, color) in enumerate(perf_data):
    left = Inches(0.7) + i * Inches(4.1)
    add_text_box(slide, left, Inches(4.4), Inches(3.8), Inches(0.4),
                 metric, 22, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(4.85), Inches(3.8), Inches(0.35),
                 detail, 14, TEXT_LIGHT, False, PP_ALIGN.CENTER)

# 未来展望
add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.35),
             "未来展望", 20, ACCENT_BLUE, True)

future_items = [
    "① 探索更多数据集验证泛化性 (如其他呼吸音数据库)",
    "② 研究与其他 backbone 的兼容性 (如 HTS-AT, Swin Transformer)",
    "③ 扩展至更多分类任务 (如2类异常/正常分类, 3类诊断分类)",
    "④ 结合更强数据增强策略进一步提升性能",
]
para_data = [(item, 15, TEXT_LIGHT, False, PP_ALIGN.LEFT) for item in future_items]
add_multi_para(slide, Inches(0.7), Inches(5.9), Inches(11), Inches(1.5), para_data)

# ================================================================
# Slide 11: 参考文献
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.5),
             "10 参考文献 & 致谢", 28, ACCENT_BLUE, True)

refs = [
    "[1] Pramono et al., Automatic adventitious respiratory sound analysis, PloS one, 2017",
    "[3] Gairola et al., RespireNet, EMBC 2021",
    "[4] Kim et al., RepAugment, arXiv 2024",
    "[8] Kim et al., SG-SCL, ICASSP 2024",
    "[9] Bae et al., Patch-Mix CL, Interspeech 2023",
    "[21] Huang et al., Adaptive frequency filters as efficient global token mixers, ICCV 2023",
    "[22] Ye et al., Differential Transformer, arXiv 2024",
    "[23] Rocha et al., ICBHI respiratory sound database, 2018",
]
para_data = [(ref, 12, TEXT_LIGHT, False, PP_ALIGN.LEFT) for ref in refs]
add_multi_para(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(2.8), para_data)

# 致谢
add_card(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(2.5),
         "致谢与资助", [
             "■ 国家自然科学基金 (No. 82260024)",
             "■ 南昌大学第一附属医院",
             "■ 南昌大学影像与视觉表示实验室",
             "",
             "■ 代码开源: https://github.com/deegy666/ADD-RSC",
             "",
             "■ 感谢: patch-mix_contrastive_learning, FunASR社区, AdaptiveFrequencyFilters"
         ], ACCENT_BLUE, 18, 14)

# 论文引用
add_text_box(slide, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
             "Dong et al., Adaptive Differential Denoising for Respiratory Sounds Classification, Interspeech 2025", 14, TEXT_DIM, False, PP_ALIGN.CENTER)

# ============ 保存PPT ============
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ADD4RSC_论文介绍_中文PPT.pptx")
prs.save(output_path)
print(f"PPT已保存至: {output_path}")
